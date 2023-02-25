import PyPDF2
import openai
import time
from pikepdf import Pdf, Name, PdfImage
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import Image

def summarise(name):
    openai.api_key = "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"
    model_engine = "text-davinci-003"
    user_prompt = "Generate a powerpoint slide header and 4 bullet points that are at most 2 sentences long for this text: "

    # Open the PDF file in read binary mode
    pdf_file_download = open(os.path.join('papers', name), 'rb')

    # Create a PDF reader object
    pdf_reader = PyPDF2.PdfReader(pdf_file_download)

    # Get the total number of pages in the PDF file
    num_pages = len(pdf_reader.pages)

    # Create a text file to write the converted text
    text_file_download = open(os.path.join('output', f'{name}.txt'), 'w')

    responseTotal = []
    for page_num in range(num_pages):
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        print("new page \n",text)
        text_file_download.write(text)
        prompt = user_prompt + text

        completion = openai.Completion.create(
            engine = model_engine,
            prompt = prompt,
            max_tokens = 1024,
            n = 1,
            stop = None,
            temperature = 0.5,
        )

        response = completion.choices[0].text
        print(response + "\n")
        responseTotal.append(response)

    # Close the files
    pdf_file_download.close()
    text_file_download.close()

    return responseTotal

def extract_images(name):
    folder_path = 'images'
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path) # delete the file
            elif os.path.isdir(file_path):
                os.rmdir(file_path) # delete the folder
        except Exception as e:
            print(f"Failed to delete {file_path}. Reason: {e}")

    old_pdf = Pdf.open(os.path.join('papers', name))
    imagenames = []
    for i in range(len(old_pdf.pages)):
        page_1 = old_pdf.pages[i]
        if (list(page_1.images.keys()) != []):
            imagenames.append((list(page_1.images.keys())))
            count = 1
            for x in page_1.images.keys():
                raw_image = page_1.images[x]
                pdf_image = PdfImage(raw_image)
                pdf_image.extract_to(fileprefix= os.path.join(folder_path, f'page{i}image{count}'))
                count+=1

    print(imagenames)


def generate_slides(infoList):
    prs = Presentation()
    print(infoList)

    for item in infoList:
        # Split the summary into bullet points
        bullet_points = item['summary'].split("\n")
        bullet_points = [bp.strip() for bp in bullet_points]
        
        # Add each object as a slide
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        slide_title = slide.shapes.title
        slide_title.text = bullet_points[0].split("\n")[0]
        slide_title.text_frame.paragraphs[0].font.size = Pt(30)
        slide_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(14, 77, 22)
        slide_title.text_frame.paragraphs[0].font.bold = True
        tf = slide.shapes.placeholders[1].text_frame

        for bullet_point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = bullet_point
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0,0,0)
            p.level = 1
        
        image_path = os.path.join('images', os.path.basename(item['image_url']))
        logo_path = os.path.join(os.getcwd(), 'msd.png')
        
        slide.shapes.add_picture(logo_path, Inches(7.22), Inches(6.22), width=Inches(2.77), height=Inches(1.17))
        slide.shapes.add_picture(image_path, Inches(0.22), Inches(6.22), width=Inches(2.77), height=Inches(1.17))

    return prs
