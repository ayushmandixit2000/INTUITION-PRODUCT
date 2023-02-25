import PyPDF2
import openai
import time
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import Image
import os




openai.api_key = "YOUR-API-KEY-HERE"
model_engine = "text-davinci-003"
user_prompt = "I want a powerpoint slide header and description from this text in 4 bullet points: "


def generate_summary(pdf_path):
    # Open the PDF file in read binary mode
    with open(pdf_path, 'rb') as pdf_file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfFileReader(pdf_file)

        # Get the total number of pages in the PDF file
        num_pages = pdf_reader.getNumPages()

        response_total = ""
        for page_num in range(num_pages):
            page = pdf_reader.getPage(page_num)
            text = page.extractText()
            prompt = user_prompt + text

            completion = openai.Completion.create(
                engine=model_engine,
                prompt=prompt,
                max_tokens=1024,
                n=1,
                stop=None,
                temperature=0.5,
            )

            response = completion.choices[0].text.strip()
            response_total += response + "\n"
            time.sleep(5)

    return response_total


def generate_slides(pptx_path):
    prs = Presentation()

    # Generate the summary
    #summary = generate_summary(pdf_path)
    summary = 'Description:Helps children with ASD develop the ability to recognize emotion from human expressions\nEnhances their ability to communicate and socialize with people\nAids in their transition into adulthood and integration into society\nAllows for greater independence and success in the workforce'

    # Split the summary into bullet points
    bullet_points = summary.split("\n")
    bullet_points = [bp.strip() for bp in bullet_points]
    
    # Add each object as a slide
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide_title = slide.shapes.title
    slide_title.text = bullet_points[0].split("\n")[0]
    slide_title.text_frame.paragraphs[0].font.size = Pt(30)
    slide_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(14, 77, 22)
    slide_title.text_frame.paragraphs[0].font.bold = True
    tf = slide.shapes.placeholders[1].text_frame

    for bullet_point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bullet_point
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0,0,0)
        p.level = 1
    image_path = "/Users/someshsahu/Downloads/MSD.png"
    img = Image.open(image_path)
    pic = slide.shapes.add_picture(image_path, Inches(7.22), Inches(6.22), width=Inches(2.77), height=Inches(1.17))

    # Save the PowerPoint file
    prs.save(pptx_path)


generate_slides('/Users/someshsahu/Downloads/dryrun5.ppt')
os.system("open /Users/someshsahu/Downloads/dryrun5.ppt")